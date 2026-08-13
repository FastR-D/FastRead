"""Tests for one-click reading-report -> .pptx projection."""

import io

from pptx import Presentation

from app.services.ppt_service import build_reading_report_pptx, PresentationService


def _sample_report() -> dict:
    return {
        "version": 1,
        "generated_at": "2026-08-13T00:00:00+00:00",
        "title": "Probing Safety Filters with Adaptive Prompts",
        "executive_summary": "This paper studies how adaptive prompts evade safety filters.",
        "key_questions": [
            {
                "question": "研究问题是什么？",
                "answer": "Whether adaptive prompts bypass filters.",
                "why_it_matters": "定义了威胁模型。",
                "verification_status": "supported",
                "evidence": [
                    {
                        "exact_quote": "We propose an adaptive prompt attack.",
                        "page_start": 2,
                        "page_end": 2,
                    }
                ],
            },
            {
                "question": "方法如何一步步完成？",
                "answer": "Iterative rewrite with feedback.",
                "why_it_matters": "降低攻击成本。",
                "verification_status": "source_only",
                "evidence": [],
            },
        ],
        "process": [
            {"step": "种子生成", "description": "采样初始提示。"},
            {"step": "反馈迭代", "description": "据过滤器响应改写。"},
        ],
        "contributions": [
            {
                "title": "自适应攻击框架",
                "description": "首个公开的自适应提示攻击流水线。",
                "evidence": [
                    {"exact_quote": "Our framework adapts online.", "page_start": 4, "page_end": 5}
                ],
            }
        ],
        "limitations": ["仅在 GPT 系列上验证。", "学术身份 Gate：预印本。"],
        "terms": [{"term": "自适应提示", "explanation": "随过滤器反馈变化的提示。"}],
        "suggested_questions": ["迁移到多模态的效果？"],
        "academic_gate": {"gate_passed": False, "label": "预印本"},
        "report_grounding_status": "partial",
    }


def test_build_pptx_returns_valid_deck_bytes():
    deck = build_reading_report_pptx(_sample_report(), {"title": "Probing Safety Filters", "authors": ["Alice", "Bob"], "venue": {"label": "arXiv preprint"}, "year": 2026})
    assert isinstance(deck, bytes) and len(deck) > 1000
    prs = Presentation(io.BytesIO(deck))
    titles = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text:
                titles.append(shape.text_frame.text)
                break
    # title, overview, 2 questions, process, contributions, limitations, terms, suggested, closing
    slide_count = len(list(prs.slides))
    assert slide_count >= 8
    body = "\n".join(titles)
    assert "关键问题 1" in body or "关键问题" in body
    assert "主要过程" in body
    assert "主要贡献" in body
    assert "第 2 页" in "\n".join(shape.text_frame.text for s in prs.slides for shape in s.shapes if shape.has_text_frame)


def test_build_pptx_minimal_report_does_not_crash():
    report = {"title": "Tiny", "key_questions": [], "process": [], "contributions": [], "limitations": []}
    deck = build_reading_report_pptx(report)
    prs = Presentation(io.BytesIO(deck))
    # title + closing at minimum
    assert len(list(prs.slides)) >= 2


def test_presentation_service_missing_report_raises(tmp_path):
    from app.repositories.note_artifacts import NoteArtifactRepository
    import pytest
    repo = NoteArtifactRepository(output_dir=tmp_path)
    service = PresentationService(repo)
    with pytest.raises(ValueError):
        service.build_for_task("does-not-exist")
