"""One-click presentation generation from a reading report.

Turns the structured ``reading_report`` produced by :mod:`reading_report_service`
into a ``.pptx`` so a reader can present a paper's key questions, process and
contributions without re-authoring slides. The deck is a faithful projection of
the report — it never invents content, and every evidence quote keeps its page
range so the deck stays auditable like the rest of FastRead.
"""

from __future__ import annotations

import io
import re
from datetime import datetime, timezone

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.repositories.note_artifacts import NoteArtifactRepository


# ---- theme -----------------------------------------------------------------

ACCENT = (31, 78, 121)        # deep blue
ACCENT_SOFT = (226, 235, 244)  # pale blue panel
INK = (33, 37, 41)            # near-black body text
MUTED = (108, 117, 125)       # secondary text
PAGE_TAG = (13, 110, 89)      # green page citation

BLANK_LAYOUT = 6  # built-in "Blank" layout in the default template


# ---- helpers ---------------------------------------------------------------


def _clean(value, limit: int = 600) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip()
    return text[:limit]


def _add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))


def _style_run(run, *, size, bold=False, color=INK, italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color)
    run.font.name = "Helvetica Neue"


def _rgb(triple):
    from pptx.dml.color import RGBColor

    return RGBColor(*triple)


def _add_paragraphs(box, paragraphs, *, base_size=14, gap=6):
    """Fill a text frame with ``[(text, opts), ...]`` paragraphs."""
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for text, opts in paragraphs:
        p = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        p.space_after = Pt(opts.get("gap", gap))
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        _style_run(
            run,
            size=opts.get("size", base_size),
            bold=opts.get("bold", False),
            color=opts.get("color", INK),
            italic=opts.get("italic", False),
        )
    return frame


def _slide_title_bar(slide, title, subtitle=None):
    """Standard slide header: a colored bar with the title."""
    bar = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9.0), Inches(0.9))
    _add_paragraphs(
        bar,
        [(_clean(title, 120), {"size": 26, "bold": True, "color": ACCENT})],
        gap=0,
    )
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9.0), Inches(0.4))
        _add_paragraphs(sub, [(_clean(subtitle, 140), {"size": 13, "color": MUTED})], gap=0)
    # thin divider
    line = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.5), Inches(9.0), Pt(2)
    )  # MSO_SHAPE.RECTANGLE = 1
    line.fill.solid()
    line.fill.fore_color.rgb = _rgb(ACCENT_SOFT)
    line.line.fill.background()


def _evidence_lines(evidence: list[dict]) -> list[tuple[str, dict]]:
    """Build (text, opts) paragraphs for a question/contribution's evidence."""
    lines: list[tuple[str, dict]] = []
    for item in (evidence or [])[:4]:
        quote = _clean(item.get("exact_quote"), 220)
        if not quote:
            continue
        page_start = item.get("page_start")
        page_end = item.get("page_end")
        page = ""
        if page_start and page_end and page_start != page_end:
            page = f"  · 第 {page_start}–{page_end} 页"
        elif page_start:
            page = f"  · 第 {page_start} 页"
        lines.append((f"“{quote}”", {"size": 12, "color": MUTED, "italic": True, "gap": 2}))
        if page:
            lines.append((page, {"size": 11, "color": PAGE_TAG, "gap": 8}))
    return lines


def _status_label(status: str) -> str:
    return {
        "supported": "证据支持",
        "refuted": "证据反驳",
        "mixed": "证据混合",
        "insufficient": "证据不足",
        "data_void": "无可用数据",
        "source_risk": "来源存疑",
        "source_only": "仅来源",
    }.get(str(status or ""), "")


# ---- deck builder ----------------------------------------------------------


def build_reading_report_pptx(report: dict, paper_meta: dict | None = None) -> bytes:
    """Render ``report`` (a normalized reading report) to pptx bytes."""
    paper_meta = paper_meta or {}
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title = _clean(report.get("title") or paper_meta.get("title"), 160) or "FastRead 学术阅读报告"
    authors = paper_meta.get("authors") or []
    if isinstance(authors, str):
        authors = [a.strip() for a in authors.split(",") if a.strip()]
    venue = paper_meta.get("venue") or {}
    venue_label = ""
    if isinstance(venue, dict):
        venue_label = _clean(venue.get("label") or venue.get("name"), 120)
    elif venue:
        venue_label = _clean(venue, 120)
    subtitle_parts = []
    if authors:
        subtitle_parts.append("、".join(authors[:4]) + (" 等" if len(authors) > 4 else ""))
    if venue_label:
        subtitle_parts.append(venue_label)
    if paper_meta.get("year"):
        subtitle_parts.append(str(paper_meta["year"]))
    subtitle = " · ".join(subtitle_parts)
    generated = report.get("generated_at") or datetime.now(timezone.utc).isoformat()
    try:
        generated = datetime.fromisoformat(generated.replace("Z", "+00:00")).strftime("%Y-%m-%d")
    except Exception:
        pass

    # 1. title slide
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    panel = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = _rgb(ACCENT)
    panel.line.fill.background()
    box = _add_text_box(slide, 0.9, 2.4, 11.5, 2.2)
    _add_paragraphs(
        box,
        [
            ("FastRead 学术阅读报告", {"size": 16, "color": (255, 255, 255), "gap": 10}),
            (title, {"size": 40, "bold": True, "color": (255, 255, 255), "gap": 12}),
            (subtitle or " ", {"size": 16, "color": (200, 220, 240)}),
        ],
    )
    foot = _add_text_box(slide, 0.9, 6.6, 11.5, 0.5)
    _add_paragraphs(foot, [(f"生成于 {generated} · 基于分页原文与联网核验证据", {"size": 11, "color": (200, 220, 240)})])

    # 2. overview / executive summary
    summary = _clean(report.get("executive_summary"), 900)
    if summary:
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
        _slide_title_bar(slide, "概览", "研究问题 · 方法主线 · 核心贡献")
        box = _add_text_box(slide, 0.6, 1.9, 12.1, 5.0)
        _add_paragraphs(box, [(summary, {"size": 17, "color": INK, "gap": 10})])

    # 3. key questions — one per slide
    questions = [q for q in (report.get("key_questions") or []) if q.get("question")]
    for idx, q in enumerate(questions, 1):
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
        status = _status_label(q.get("verification_status"))
        _slide_title_bar(slide, f"关键问题 {idx}", _clean(q.get("question"), 120))
        box = _add_text_box(slide, 0.6, 1.9, 12.1, 5.2)
        paras: list[tuple[str, dict]] = []
        answer = _clean(q.get("answer"), 700)
        if answer:
            paras.append((answer, {"size": 16, "color": INK, "gap": 8}))
        why = _clean(q.get("why_it_matters"), 240)
        if why:
            paras.append((f"为什么重要：{why}", {"size": 13, "color": MUTED, "italic": True, "gap": 10}))
        if status:
            paras.append((f"证据状态：{status}", {"size": 12, "color": PAGE_TAG, "gap": 10}))
        paras.extend(_evidence_lines(q.get("evidence") or []))
        if not paras:
            paras.append(("（报告未给出正文）", {"size": 14, "color": MUTED, "italic": True}))
        _add_paragraphs(box, paras)

    # 4. main process
    process = [p for p in (report.get("process") or []) if p.get("step") or p.get("description")]
    if process:
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
        _slide_title_bar(slide, "主要过程", "方法如何一步步完成")
        box = _add_text_box(slide, 0.6, 1.9, 12.1, 5.2)
        paras = []
        for idx, step in enumerate(process, 1):
            label = _clean(step.get("step") or f"步骤 {idx}", 80)
            desc = _clean(step.get("description"), 320)
            paras.append((f"{idx}. {label}", {"size": 16, "bold": True, "color": ACCENT, "gap": 2}))
            if desc:
                paras.append((desc, {"size": 14, "color": INK, "gap": 10}))
        _add_paragraphs(box, paras)

    # 5. contributions
    contribs = [c for c in (report.get("contributions") or []) if c.get("title") or c.get("description")]
    if contribs:
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
        _slide_title_bar(slide, "主要贡献", "相对已有工作的增量")
        box = _add_text_box(slide, 0.6, 1.9, 12.1, 5.2)
        paras = []
        for idx, c in enumerate(contribs, 1):
            title_c = _clean(c.get("title") or f"贡献 {idx}", 100)
            desc = _clean(c.get("description"), 360)
            paras.append((f"{idx}. {title_c}", {"size": 16, "bold": True, "color": ACCENT, "gap": 2}))
            if desc:
                paras.append((desc, {"size": 14, "color": INK, "gap": 6}))
            paras.extend(_evidence_lines(c.get("evidence") or []))
        _add_paragraphs(box, paras)

    # 6. limitations
    limitations = [str(l).strip() for l in (report.get("limitations") or []) if str(l).strip()]
    if limitations:
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
        _slide_title_bar(slide, "局限与证据边界", "报告未能覆盖或需谨慎之处")
        box = _add_text_box(slide, 0.6, 1.9, 12.1, 5.2)
        paras = [(f"• { _clean(l, 360)}", {"size": 15, "color": INK, "gap": 8}) for l in limitations]
        _add_paragraphs(box, paras)

    # 7. terms
    terms = [t for t in (report.get("terms") or []) if t.get("term")]
    if terms:
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
        _slide_title_bar(slide, "关键术语")
        box = _add_text_box(slide, 0.6, 1.9, 12.1, 5.2)
        paras = []
        for t in terms:
            paras.append((_clean(t.get("term"), 60), {"size": 15, "bold": True, "color": ACCENT, "gap": 1}))
            paras.append((_clean(t.get("explanation"), 300), {"size": 13, "color": INK, "gap": 8}))
        _add_paragraphs(box, paras)

    # 8. suggested follow-up questions
    suggested = [str(s).strip() for s in (report.get("suggested_questions") or []) if str(s).strip()]
    if suggested:
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
        _slide_title_bar(slide, "可继续追问", "带着这些问题回到原文")
        box = _add_text_box(slide, 0.6, 1.9, 12.1, 5.2)
        paras = [(f"{idx}. { _clean(s, 200)}", {"size": 15, "color": INK, "gap": 8}) for idx, s in enumerate(suggested, 1)]
        _add_paragraphs(box, paras)

    # 9. closing
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT])
    panel = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = _rgb(ACCENT)
    panel.line.fill.background()
    box = _add_text_box(slide, 1.0, 3.0, 11.3, 1.5)
    _add_paragraphs(
        box,
        [
            ("感谢阅读", {"size": 34, "bold": True, "color": (255, 255, 255), "gap": 8}),
            ("本幻灯片由 FastRead 阅读报告自动生成，请结合分页原文核对引文。", {"size": 14, "color": (200, 220, 240)}),
        ],
    )

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


class PresentationService:
    """Reads a task's reading report + paper metadata and builds a deck."""

    def __init__(self, artifacts: NoteArtifactRepository | None = None):
        self.artifacts = artifacts or NoteArtifactRepository()

    def build_for_task(self, task_id: str) -> tuple[bytes, str]:
        result = self.artifacts.read_result(task_id)
        if not result:
            raise ValueError("任务不存在或尚未导入正文")
        report = (result.get("insights") or {}).get("reading_report")
        if not report:
            raise ValueError("尚未生成阅读报告，请先一键生成关键问题报告")
        paper_doc = result.get("paper_document") or {}
        audio_meta = result.get("audio_meta") or {}
        paper_meta = {
            "title": paper_doc.get("title") or audio_meta.get("title") or report.get("title"),
            "authors": paper_doc.get("authors") or (audio_meta.get("raw_info") or {}).get("authors"),
            "venue": paper_doc.get("venue"),
            "year": paper_doc.get("year"),
        }
        deck = build_reading_report_pptx(report, paper_meta)
        base = _clean(paper_meta.get("title") or "reading-report", 60) or "reading-report"
        safe = re.sub(r"[^\w一-鿿\-]+", "_", base).strip("_") or "reading-report"
        filename = f"{safe}.pptx"
        return deck, filename
